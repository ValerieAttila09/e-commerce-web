#!/usr/bin/env python3
# Simple Markdown -> PPTX converter for this repo's presentation markdown
# Requires: python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
import re
import sys

MD_PATH = 'docs/presentation_5min_slides.md'
OUT_PATH = 'docs/presentation_5min_slides.pptx'

with open(MD_PATH, 'r', encoding='utf-8') as f:
    md = f.read()

# Split into sections starting with '## '
parts = re.split(r'^##\s+', md, flags=re.MULTILINE)
# first part is preface (skip until first slide)
parts = [p.strip() for p in parts if p.strip()]

prs = Presentation()
# set slide dimensions / default font sizes via placeholders

for part in parts:
    # The part starts with the rest of the header line (title)
    # Extract title (first line) and the rest
    lines = part.splitlines()
    title = lines[0].strip()
    body = '\n'.join(lines[1:]).strip()

    # Create slide with title and content
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Find the 'Tampilan slide' list and 'Catatan pembicara'
    # We'll try to extract the bullet points that appear under 'Tampilan slide (poin singkat):' or 'Tampilan slide (poin singkat):'
    bullets = []
    notes = ''

    # Try to extract section between 'Tampilan slide' and 'Catatan pembicara'
    m = re.search(r"Tampilan slide.*?:\s*(.*?)\n\n", body, flags=re.S)
    if m:
        candidate = m.group(1)
        # split lines and take lines starting with '-'
        for line in candidate.splitlines():
            line = line.strip()
            if line.startswith('- '):
                bullets.append(line[2:].strip())

    # Alternative: look for lines like '- Judul: ...' earlier
    if not bullets:
        for ln in body.splitlines():
            ln = ln.strip()
            if ln.startswith('- '):
                bullets.append(ln[2:].strip())

    # If still empty, take the first paragraph as bullet
    if not bullets:
        # take first 4 short lines as bullets
        for ln in body.splitlines():
            ln = ln.strip()
            if ln:
                bullets.append(ln)
            if len(bullets) >= 4:
                break

    # Populate content placeholder as bullet points
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
            p.level = 0
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    # Extract notes under 'Catatan pembicara' until next '---' or end
    m2 = re.search(r"Catatan pembicara.*?:\s*(.*)", body, flags=re.S)
    if m2:
        notes = m2.group(1).strip()
        # Clean trailing separators
        notes = re.split(r'---|\n## |\n\n##', notes)[0].strip()
    else:
        # As fallback, try to extract 'Catatan pembicara (skrip' or 'Catatan pembicara' section
        m3 = re.search(r"Catatan pembicara.*?\)?:\s*(.*?)$", body, flags=re.S)
        if m3:
            notes = m3.group(1).strip()

    # Add notes to slide
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

# Save presentation
prs.save(OUT_PATH)
print('Saved:', OUT_PATH)
